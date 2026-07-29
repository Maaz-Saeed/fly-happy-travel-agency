"""Generates the class presentation deck for the Fly Happy International Travels project.

Run with:  python scripts/generate_presentation.py
Output:    Fly_Happy_Presentation.pptx (project root)
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
IMG = os.path.join(BASE, "app", "static", "images")
OUT = os.path.join(BASE, "Fly_Happy_Presentation.pptx")

# ---- Brand palette ----
DARK = RGBColor(0x0D, 0x4F, 0x31)
GREEN = RGBColor(0x14, 0x7A, 0x49)
LIGHT = RGBColor(0x2E, 0xA8, 0x6B)
PALE = RGBColor(0xE8, 0xF5, 0xEE)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARKTEXT = RGBColor(0x1C, 0x2B, 0x24)
MUTED = RGBColor(0x5F, 0x72, 0x68)

SW, SH = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=DARKTEXT, bold=False, align=PP_ALIGN.LEFT,
             font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=DARKTEXT, bullet_color=GREEN, gap=6, bold_lead=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = "▸  " + lead
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = GREEN; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = "  " + rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "▸  " + item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return box


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.15), DARK)
    add_rect(slide, 0, Inches(1.15), SW, Pt(4), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.6), title,
              size=28, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.72), Inches(11.5), Inches(0.4), subtitle,
                  size=14, color=GOLD, italic=True)
    # small logo top-right
    try:
        slide.shapes.add_picture(os.path.join(IMG, "logo.png"), Inches(12.35), Inches(0.15), height=Inches(0.85))
    except Exception:
        pass


def footer(slide, page_no):
    add_rect(slide, 0, SH - Inches(0.4), SW, Inches(0.4), PALE)
    add_text(slide, Inches(0.4), SH - Inches(0.38), Inches(8), Inches(0.35),
              "Fly Happy International Travels  |  License No. PR-5199", size=10, color=MUTED)
    add_text(slide, SW - Inches(1.4), SH - Inches(0.38), Inches(1.0), Inches(0.35),
              str(page_no), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def new_slide(title=None, subtitle=None, page_no=None):
    slide = prs.slides.add_slide(BLANK)
    set_background(slide, WHITE)
    if title:
        header(slide, title, subtitle)
    if page_no:
        footer(slide, page_no)
    return slide


def stat_card(slide, x, y, w, h, number, label, color):
    add_rect(slide, x, y, w, h, color)
    add_text(slide, x, y + Inches(0.12), w, Inches(0.55), number, size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + h - Inches(0.5), w, Inches(0.4), label, size=11, color=WHITE, align=PP_ALIGN.CENTER)


def box_node(slide, x, y, w, h, text, fill=GREEN, txt_color=WHITE, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = DARK; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txt_color; r.font.name = "Calibri"
    return shp


def arrow(slide, x, y, w, h, rotation=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.rotation = rotation
    shp.shadow.inherit = False
    return shp


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
slide = prs.slides.add_slide(BLANK)
set_background(slide, DARK)
add_rect(slide, 0, Inches(6.6), SW, Inches(0.9), GREEN)
add_rect(slide, 0, Inches(6.55), SW, Pt(4), GOLD)
try:
    slide.shapes.add_picture(os.path.join(IMG, "logo.png"), Inches(5.66), Inches(0.6), height=Inches(2.0))
except Exception:
    pass
add_text(slide, Inches(1), Inches(2.85), Inches(11.33), Inches(1.1),
          "FLY HAPPY INTERNATIONAL TRAVELS", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11.33), Inches(0.6),
          "Travel Agency Management System", size=22, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, Inches(1), Inches(4.35), Inches(11.33), Inches(0.5),
          "A Full-Stack Python Flask Web Application", size=15, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.72), Inches(11.33), Inches(0.6),
          "Presented by Maaz Saeed  |  Reg. No: SU-23-01-002-031  |  BS Software Engineering, 6th Semester",
          size=13, color=WHITE, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 — AGENDA
# =====================================================================
slide = new_slide("Presentation Agenda", page_no=2)
agenda = [
    "Project Brief — What Was Requested",
    "Objectives",
    "Technology Stack",
    "System Architecture",
    "Database Design",
    "Public Website Walkthrough",
    "Booking & Payment Flow",
    "Customer Dashboard",
    "Admin Panel & Analytics",
    "Chatbot Assistant",
    "Security Measures",
    "Testing, Verification & Documentation",
    "Future Enhancements",
    "Conclusion",
]
col1 = agenda[:7]
col2 = agenda[7:]
add_bullets(slide, Inches(0.7), Inches(1.5), Inches(5.9), Inches(5.3), col1, size=17, gap=14)
add_bullets(slide, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.3), col2, size=17, gap=14)

# =====================================================================
# SLIDE 3 — PROJECT BRIEF (the client prompt)
# =====================================================================
slide = new_slide("Project Brief", "What the client asked for", page_no=3)
add_text(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.8),
          "Build a COMPLETE, PROFESSIONAL, CLIENT-READY Travel Agency Management System — "
          "not a student-level project, but something that could be sold to a real travel agency.",
          size=16, color=DARKTEXT, italic=True)
brief_items = [
    ("Company:", "Fly Happy International Travels — License No. PR-5199, Green & White theme"),
    ("Stack:", "Python Flask, HTML5/CSS3/Bootstrap 5/JavaScript, SQLite, Jinja2, Flask Sessions"),
    ("Scope requested:", "Home/About/Services pages, 80+ destinations, 40+ airlines, booking + payment "
                          "simulation, customer dashboard, full admin panel, chatbot, reports, security, and "
                          "complete SRS/ER/manuals/viva documentation"),
    ("Delivery style:", "One working project — no placeholders, no \"continue later\", fully runnable locally"),
]
add_bullets(slide, Inches(0.7), Inches(2.5), Inches(11.9), Inches(4.3), brief_items, size=16, gap=16)

# =====================================================================
# SLIDE 4 — OBJECTIVES
# =====================================================================
slide = new_slide("Project Objectives", page_no=4)
objs = [
    "Deliver an end-to-end booking journey: browse → book → pay → manage",
    "Model a realistic, normalized database covering the full agency operation",
    "Give staff a complete back-office: catalog, bookings, payments, staff, content, analytics",
    "Bake in security from the start — hashed passwords, CSRF, ORM-only queries, access control",
    "Make the UI feel like a real commercial product (animation, branding, responsiveness)",
    "Produce professional documentation suitable for handover and academic evaluation",
]
add_bullets(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5), objs, size=18, gap=22)

# =====================================================================
# SLIDE 5 — TECH STACK
# =====================================================================
slide = new_slide("Technology Stack", page_no=5)
stack = [
    ("Backend", "Python 3, Flask 3 (Application Factory + Blueprints)"),
    ("Frontend", "HTML5, CSS3, Bootstrap 5, Vanilla JavaScript"),
    ("Database", "SQLite + SQLAlchemy ORM (Flask-SQLAlchemy)"),
    ("Templating", "Jinja2 (template inheritance: base.html / admin/_base.html)"),
    ("Auth", "Flask-Login (sessions) + Werkzeug password hashing"),
    ("Forms & CSRF", "Flask-WTF / WTForms"),
    ("Icons", "Font Awesome 6"),
    ("Charts", "Chart.js (dashboard analytics)"),
    ("Animation", "AOS (Animate On Scroll) + custom CSS/JS animations"),
]
rows, cols = 3, 3
cw, ch = Inches(3.9), Inches(1.55)
for i, (k, v) in enumerate(stack):
    r, c = divmod(i, cols)
    x = Inches(0.6) + c * (cw + Inches(0.25))
    y = Inches(1.5) + r * (ch + Inches(0.25))
    add_rect(slide, x, y, cw, ch, PALE)
    add_rect(slide, x, y, Inches(0.12), ch, GREEN)
    add_text(slide, x + Inches(0.3), y + Inches(0.12), cw - Inches(0.5), Inches(0.4), k, size=16, bold=True, color=DARK)
    add_text(slide, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.5), Inches(0.9), v, size=12.5, color=DARKTEXT)

# =====================================================================
# SLIDE 6 — SYSTEM ARCHITECTURE
# =====================================================================
slide = new_slide("System Architecture", "Flask Application Factory + Blueprint pattern", page_no=6)
box_node(slide, Inches(5.1), Inches(1.5), Inches(3.1), Inches(0.7), "Web Browser\n(Bootstrap, JS, Chart.js)", fill=DARK)
arrow(slide, Inches(6.4), Inches(2.25), Inches(0.5), Inches(0.35), rotation=90)
box_node(slide, Inches(3.2), Inches(2.75), Inches(6.9), Inches(0.8),
         "Flask Blueprints (Controllers): public / auth / booking / payment / customer / admin / chatbot", fill=GREEN)
arrow(slide, Inches(6.4), Inches(3.6), Inches(0.5), Inches(0.35), rotation=90)
box_node(slide, Inches(3.9), Inches(4.0), Inches(5.5), Inches(0.7), "Forms Layer — Flask-WTF (validation + CSRF)", fill=LIGHT, txt_color=DARK)
arrow(slide, Inches(6.4), Inches(4.75), Inches(0.5), Inches(0.35), rotation=90)
box_node(slide, Inches(3.9), Inches(5.15), Inches(5.5), Inches(0.7), "Models — SQLAlchemy ORM (16 tables)", fill=GREEN)
arrow(slide, Inches(6.4), Inches(5.9), Inches(0.5), Inches(0.35), rotation=90)
box_node(slide, Inches(4.6), Inches(6.3), Inches(4.1), Inches(0.65), "SQLite Database\ninstance/flyhappy.db", fill=DARK)

# =====================================================================
# SLIDE 7 — DATABASE DESIGN
# =====================================================================
slide = new_slide("Database Design", "16 normalized tables (3NF)", page_no=7)
tables = ["users", "employees", "destinations", "airlines", "packages", "hotels",
          "bookings", "payments", "testimonials", "team_members", "news", "faqs",
          "gallery", "contact_messages", "chat_logs", "newsletter_subscribers"]
cols = 4
cw, ch = Inches(2.9), Inches(0.65)
for i, t in enumerate(tables):
    r, c = divmod(i, cols)
    x = Inches(0.65) + c * (cw + Inches(0.2))
    y = Inches(1.55) + r * (ch + Inches(0.2))
    box_node(slide, x, y, cw, ch, t, fill=GREEN if (i % 2 == 0) else LIGHT, size=13)
add_text(slide, Inches(0.65), Inches(4.55), Inches(12), Inches(1.9),
          "Key relationships:\n"
          "•  users (1) —< bookings >— (1) destinations, bookings >— (1) airlines\n"
          "•  bookings (1) —< payments   |   users (1) — (1) employees\n"
          "•  destinations (1) —< packages, hotels   |   users (1) —< chat_logs",
          size=15, color=DARKTEXT)

# =====================================================================
# SLIDE 8 — PUBLIC WEBSITE
# =====================================================================
slide = new_slide("Public Website", "Marketing site with 85 destinations & 28 airlines", page_no=8)
pages = [
    "Home — hero, animated stats, featured destinations/packages/hotels, testimonials, news, gallery",
    "About — mission, vision, core values, timeline, CEO message, awards",
    "Services — 26 service cards (ticketing, visas, Hajj/Umrah, insurance, and more)",
    "Destinations — 85 destinations with search, region & visa filters, pagination",
    "Airlines, Packages, Hotels, Gallery, Team, News, FAQ, Testimonials, Contact (with map)",
]
add_bullets(slide, Inches(0.7), Inches(1.5), Inches(7.1), Inches(5), pages, size=15.5, gap=18)
try:
    slide.shapes.add_picture(os.path.join(IMG, "destinations", "dxb.jpg"), Inches(8.1), Inches(1.55), width=Inches(2.2))
    slide.shapes.add_picture(os.path.join(IMG, "destinations", "ist.jpg"), Inches(10.45), Inches(1.55), width=Inches(2.2))
    slide.shapes.add_picture(os.path.join(IMG, "destinations", "mle.jpg"), Inches(8.1), Inches(3.85), width=Inches(2.2))
    slide.shapes.add_picture(os.path.join(IMG, "destinations", "lhr.jpg"), Inches(10.45), Inches(3.85), width=Inches(2.2))
except Exception:
    pass

# =====================================================================
# SLIDE 9 — BOOKING & PAYMENT FLOW
# =====================================================================
slide = new_slide("Booking & Payment Flow", page_no=9)
steps = ["Search\nDestination", "Fill Trip &\nPassenger Details", "Auto-Calculated\nPrice", "Booking ID +\nTicket Number", "Choose Payment\nMethod", "Confirmed +\nInvoice"]
x = Inches(0.4)
w = Inches(1.95)
for i, s in enumerate(steps):
    box_node(slide, x, Inches(2.0), w, Inches(1.3), s, fill=GREEN if i % 2 == 0 else DARK, size=12.5)
    if i < len(steps) - 1:
        arrow(slide, x + w, Inches(2.5), Inches(0.25), Inches(0.3))
    x += w + Inches(0.25)
add_bullets(slide, Inches(0.9), Inches(4.1), Inches(11.3), Inches(2.7), [
    "Price = base fare × seat-class multiplier (Economy 1.0 / Business 1.6 / First 2.2) × "
    "trip-type multiplier (One Way 1.0 / Round Trip 1.85) × passengers (infants at 25%)",
    "Payment methods: Visa Card, MasterCard, Debit Card, JazzCash, EasyPaisa, Bank Transfer, Cash",
    "Successful payment → booking auto-confirms; Cash → stays pending until office visit",
    "Printable receipt and full itemized invoice generated with company logo & license number",
], size=15, gap=12)

# =====================================================================
# SLIDE 10 — CUSTOMER DASHBOARD
# =====================================================================
slide = new_slide("Customer Dashboard", page_no=10)
cust = [
    "View all bookings, filter by status (Pending / Confirmed / Completed / Cancelled)",
    "Cancel an active booking directly from the dashboard",
    "Download / print ticket confirmation and invoice",
    "View complete payment history with receipts",
    "Update profile details (name, phone, CNIC, passport, address)",
    "Change password securely (current password re-verified)",
    "Submit a support request straight to the admin team",
]
add_bullets(slide, Inches(0.9), Inches(1.6), Inches(11.3), Inches(5.3), cust, size=17, gap=18)

# =====================================================================
# SLIDE 11 — ADMIN PANEL
# =====================================================================
slide = new_slide("Admin Panel & Analytics", page_no=11)
stat_card(slide, Inches(0.6), Inches(1.5), Inches(2.85), Inches(1.1), "6", "Bookings Today", GREEN)
stat_card(slide, Inches(3.6), Inches(1.5), Inches(2.85), Inches(1.1), "PKR 2.4M", "Total Revenue", DARK)
stat_card(slide, Inches(6.6), Inches(1.5), Inches(2.85), Inches(1.1), "85", "Destinations", LIGHT)
stat_card(slide, Inches(9.6), Inches(1.5), Inches(2.85), Inches(1.1), "28", "Airlines", GOLD)
admin_items = [
    "Dashboard: revenue / bookings / customer-growth charts (Chart.js), live KPIs",
    "Full CRUD: destinations, airlines, packages, hotels, customers, employees",
    "Bookings & Payments: search, filter, status updates, CSV export, printable reports",
    "Content management: gallery, testimonials (approve/hide), news, FAQs",
    "Support: contact messages inbox, chatbot conversation logs",
]
add_bullets(slide, Inches(0.7), Inches(2.9), Inches(11.9), Inches(3.9), admin_items, size=16, gap=14)

# =====================================================================
# SLIDE 12 — CHATBOT
# =====================================================================
slide = new_slide("Chatbot Assistant", "Rule-based virtual travel assistant", page_no=12)
chat_items = [
    "Answers: office timings, visa process, passport requirements, booking steps, refund policy, "
    "payment methods, insurance, destinations, office location/phone/email, flight status guidance, packages",
    "Keyword-matching engine — transparent, predictable, zero external API cost",
    "Every conversation (session, message, response, timestamp) is logged and reviewable by admins",
    "Floating widget available on every public page",
]
add_bullets(slide, Inches(0.9), Inches(1.7), Inches(11.3), Inches(4.5), chat_items, size=17, gap=20)

# =====================================================================
# SLIDE 13 — SECURITY
# =====================================================================
slide = new_slide("Security Measures", page_no=13)
sec_items = [
    ("Password Hashing:", "Werkzeug PBKDF2-SHA256 — plaintext passwords are never stored"),
    ("CSRF Protection:", "Flask-WTF CSRFProtect on every state-changing form (WTForms + plain HTML)"),
    ("SQL Injection Prevention:", "100% SQLAlchemy ORM query builder — no raw SQL string interpolation"),
    ("Access Control:", "role-based (customer/employee/admin) + row-level ownership checks on bookings/payments"),
    ("Session Hardening:", "HttpOnly + SameSite cookies, bounded session lifetime"),
    ("Error Handling:", "Custom 403 / 404 / 500 pages, DB rollback on server error"),
]
add_bullets(slide, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.3), sec_items, size=16, gap=18)

# =====================================================================
# SLIDE 14 — TESTING, VERIFICATION & DOCUMENTATION
# =====================================================================
slide = new_slide("Testing, Verification & Documentation", page_no=14)
add_text(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4), "Verified live end-to-end:", size=17, bold=True, color=DARK)
verify_items = [
    "All 17 public routes → HTTP 200",
    "Customer login → dashboard, Admin login → panel (all sub-pages HTTP 200)",
    "Full flow: booking → JazzCash payment → confirmed → invoice generated correctly",
    "Chatbot API replies verified for multiple query types",
    "CSV export, admin add/edit/delete, and custom 404 page all confirmed working",
]
add_bullets(slide, Inches(0.9), Inches(1.9), Inches(11.3), Inches(2.4), verify_items, size=15, gap=10)
add_text(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4), "Documentation suite (docs/):", size=17, bold=True, color=DARK)
docs_items = ["SRS", "Database Design", "ER Diagram", "Use Cases", "System Architecture",
              "Installation Guide", "User Manual", "Admin Manual", "Testing Plan",
              "Future Enhancements", "Viva Q&A"]
cols = 4
cw, ch = Inches(2.9), Inches(0.55)
for i, d in enumerate(docs_items):
    r, c = divmod(i, cols)
    x = Inches(0.65) + c * (cw + Inches(0.2))
    y = Inches(5.0) + r * (ch + Inches(0.15))
    box_node(slide, x, y, cw, ch, d, fill=LIGHT, txt_color=DARK, size=12)

# =====================================================================
# SLIDE 15 — FUTURE ENHANCEMENTS
# =====================================================================
slide = new_slide("Future Enhancements", page_no=15)
future_items = [
    "Integrate a real payment gateway (Stripe / PayFast / JazzCash / EasyPaisa live APIs)",
    "Email/SMS notifications for booking confirmations and reminders",
    "True PDF export for invoices and tickets",
    "File-upload widgets for admin content management",
    "Migrate SQLite → PostgreSQL for multi-branch, high-concurrency scale",
    "AI-powered chatbot upgrade (LLM-backed, live booking lookups)",
    "Automated pytest test suite + CI pipeline",
]
add_bullets(slide, Inches(0.9), Inches(1.6), Inches(11.3), Inches(5.3), future_items, size=16.5, gap=16)

# =====================================================================
# SLIDE 16 — CONCLUSION / THANK YOU
# =====================================================================
slide = prs.slides.add_slide(BLANK)
set_background(slide, DARK)
try:
    slide.shapes.add_picture(os.path.join(IMG, "logo.png"), Inches(5.66), Inches(0.9), height=Inches(1.7))
except Exception:
    pass
add_text(slide, Inches(1), Inches(2.9), Inches(11.33), Inches(0.8), "Thank You", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.75), Inches(11.33), Inches(0.6),
          "A complete, secure, full-stack travel agency platform — built with Python Flask",
          size=16, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, Inches(1), Inches(4.6), Inches(11.33), Inches(0.5), "Questions?", size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
          "Maaz Saeed  |  SU-23-01-002-031  |  BS Software Engineering", size=13, color=WHITE, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(list(prs.slides)))
